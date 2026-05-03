#!/usr/bin/env python3
"""Document Parser MCP Server — extract structured content from common document formats.

Parses .pdf, .xlsx, .xlsm, .csv, .docx, and .pptx files into JSON or Markdown.
This is a read-only parser/extractor: it never writes or generates document files,
and it never executes macros from .xlsm files.

Exposes MCP tools:

* ``parse_document``  — auto-detect format by extension and parse to JSON or Markdown.
* ``parse_pdf``       — parse a PDF file.
* ``parse_excel``     — parse an Excel (.xlsx/.xlsm) file (read-only, no macro execution).
* ``parse_csv``       — parse a CSV file.
* ``parse_docx``      — parse a Word (.docx) file.
* ``parse_pptx``      — parse a PowerPoint (.pptx) file.

Run with:
    python3 server.py          # stdio transport (default)
"""

from __future__ import annotations

import csv
import io
import json
import os
import sys
import time
from pathlib import Path
from typing import Any

from mcp.server.fastmcp import FastMCP

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB — reject files larger than this
SUPPORTED_EXTENSIONS = {".pdf", ".xlsx", ".xlsm", ".csv", ".docx", ".pptx"}
OUTPUT_FORMATS = ("json", "markdown")
# CSV sniff sample size
CSV_SNIFF_BYTES = 8192
# Maximum rows to parse from a spreadsheet/CSV to prevent unbounded output
MAX_ROWS = 100_000
# Maximum characters in final tool output
MAX_OUTPUT_CHARS = 200_000

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------


def _log(msg: str) -> None:
    """Write a timestamped line to stderr (visible in MCP server logs)."""
    ts = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    print(f"[document-parser {ts}] {msg}", file=sys.stderr, flush=True)


# ---------------------------------------------------------------------------
# Error helpers
# ---------------------------------------------------------------------------


def _make_error(code: str, message: str, **extra: Any) -> dict[str, Any]:
    """Build a structured error envelope consistent with repo conventions."""
    result: dict[str, Any] = {"success": False, "error_code": code, "error": message}
    result.update(extra)
    return result


def _make_dependency_error(code: str, dependency: str, exc: Exception, file_path: Path) -> dict[str, Any]:
    """Return a structured dependency/import failure."""
    return _make_error(
        code,
        f"Required dependency '{dependency}' is unavailable: {exc}",
        dependency=dependency,
        file=str(file_path),
    )


def _validate_file(file_path: str, allowed_extensions: set[str] | None = None) -> tuple[Path, str | None]:
    """Validate that file_path exists, is within size limits, and has an allowed extension.

    Returns (resolved_path, error_message).  error_message is None on success.
    """
    p = Path(file_path).expanduser().resolve()
    if not p.exists():
        return p, f"File not found: {p}"
    if not p.is_file():
        return p, f"Path is not a file: {p}"
    ext = p.suffix.lower()
    exts = allowed_extensions or SUPPORTED_EXTENSIONS
    if ext not in exts:
        return p, f"Unsupported file extension '{ext}'. Supported: {sorted(exts)}"
    try:
        size = p.stat().st_size
    except OSError as exc:
        return p, f"Cannot stat file: {exc}"
    if size > MAX_FILE_SIZE:
        return p, f"File too large ({size:,} bytes). Maximum: {MAX_FILE_SIZE:,} bytes."
    if size == 0:
        return p, "File is empty (0 bytes)."
    return p, None


def _truncate_output(text: str) -> tuple[str, bool]:
    """Truncate text to MAX_OUTPUT_CHARS.  Returns (text, was_truncated)."""
    if len(text) <= MAX_OUTPUT_CHARS:
        return text, False
    return text[:MAX_OUTPUT_CHARS] + "\n\n… [output truncated]", True


def _json_size(value: Any) -> int:
    """Return the JSON-serialised size of a result payload."""
    return len(json.dumps(value, default=str))


def _truncate_string_value(value: Any, max_chars: int) -> Any:
    """Clip long string values while leaving non-strings untouched."""
    if not isinstance(value, str) or len(value) <= max_chars:
        return value
    return value[:max_chars] + "… [truncated]"


def _clip_table_rows(rows: list[list[Any]], max_rows: int, max_cell_chars: int = 200) -> list[list[Any]]:
    """Return a preview slice of tabular rows with clipped string cells."""
    clipped: list[list[Any]] = []
    for row in rows[:max_rows]:
        clipped.append([_truncate_string_value(cell, max_cell_chars) for cell in row])
    return clipped


def _build_truncated_json_result(parsed: dict[str, Any], original_size: int) -> dict[str, Any]:
    """Reduce oversized JSON results to a bounded preview payload."""
    fmt = parsed.get("format")
    source_truncated = _source_data_truncated(parsed)
    preview: dict[str, Any] = {
        "success": True,
        "format": fmt,
        "file": parsed.get("file"),
        "output_format": "json",
        "truncated": True,
        "output_truncated": True,
        "source_truncated": source_truncated,
        "_warning": f"Output truncated to {MAX_OUTPUT_CHARS:,} chars from {original_size:,}.",
    }

    if fmt == "csv":
        rows = parsed.get("rows", [])
        preview["delimiter"] = parsed.get("delimiter", ",")
        preview["row_count"] = parsed.get("row_count", len(rows))
        preview_rows = _clip_table_rows(rows, max_rows=200)
        preview["rows"] = preview_rows
        omitted_rows = max(0, len(rows) - len(preview_rows))
        if omitted_rows:
            preview["omitted_rows"] = omitted_rows
    elif fmt == "excel":
        sheets = parsed.get("sheets", [])
        preview["sheet_count"] = parsed.get("sheet_count", len(sheets))
        preview_sheets: list[dict[str, Any]] = []
        for sheet in sheets[:5]:
            sheet_rows = sheet.get("rows", [])
            preview_sheet = {
                "sheet_name": sheet.get("sheet_name"),
                "row_count": sheet.get("row_count", len(sheet_rows)),
                "truncated": bool(sheet.get("truncated")),
                "rows": _clip_table_rows(sheet_rows, max_rows=50),
            }
            omitted_rows = max(0, len(sheet_rows) - len(preview_sheet["rows"]))
            if omitted_rows:
                preview_sheet["omitted_rows"] = omitted_rows
            preview_sheets.append(preview_sheet)
        preview["sheets"] = preview_sheets
        omitted_sheets = max(0, len(sheets) - len(preview_sheets))
        if omitted_sheets:
            preview["omitted_sheets"] = omitted_sheets
    elif fmt == "pdf":
        pages = parsed.get("pages", [])
        preview["page_count"] = parsed.get("page_count", len(pages))
        preview["metadata"] = parsed.get("metadata", {})
        preview_pages: list[dict[str, Any]] = []
        for page in pages[:5]:
            page_preview = {"page_number": page.get("page_number")}
            if "text" in page:
                page_preview["text"] = _truncate_string_value(page.get("text", ""), 4_000)
            if page.get("tables"):
                page_preview["tables"] = [_clip_table_rows(table, max_rows=20) for table in page["tables"][:3]]
                omitted_tables = max(0, len(page["tables"]) - len(page_preview["tables"]))
                if omitted_tables:
                    page_preview["omitted_tables"] = omitted_tables
            preview_pages.append(page_preview)
        preview["pages"] = preview_pages
        omitted_pages = max(0, len(pages) - len(preview_pages))
        if omitted_pages:
            preview["omitted_pages"] = omitted_pages
    elif fmt == "docx":
        paragraphs = parsed.get("paragraphs", [])
        tables = parsed.get("tables", [])
        preview["metadata"] = parsed.get("metadata", {})
        preview["paragraph_count"] = parsed.get("paragraph_count", len(paragraphs))
        preview["table_count"] = parsed.get("table_count", len(tables))
        preview["paragraphs"] = [
            {
                "text": _truncate_string_value(paragraph.get("text", ""), 1_000),
                "style": paragraph.get("style"),
            }
            for paragraph in paragraphs[:100]
        ]
        preview["tables"] = [_clip_table_rows(table, max_rows=20) for table in tables[:5]]
        omitted_paragraphs = max(0, len(paragraphs) - len(preview["paragraphs"]))
        omitted_tables = max(0, len(tables) - len(preview["tables"]))
        if omitted_paragraphs:
            preview["omitted_paragraphs"] = omitted_paragraphs
        if omitted_tables:
            preview["omitted_tables"] = omitted_tables
    elif fmt == "pptx":
        slides = parsed.get("slides", [])
        preview["slide_count"] = parsed.get("slide_count", len(slides))
        preview_slides: list[dict[str, Any]] = []
        for slide in slides[:10]:
            preview_shapes: list[dict[str, Any]] = []
            for shape in slide.get("shapes", [])[:20]:
                shape_preview = {
                    "shape_id": shape.get("shape_id"),
                    "name": shape.get("name"),
                    "shape_type": shape.get("shape_type"),
                }
                if "text" in shape:
                    shape_preview["text"] = _truncate_string_value(shape.get("text", ""), 1_000)
                if "table" in shape:
                    shape_preview["table"] = _clip_table_rows(shape.get("table", []), max_rows=20)
                preview_shapes.append(shape_preview)
            slide_preview = {
                "slide_number": slide.get("slide_number"),
                "shapes": preview_shapes,
            }
            omitted_shapes = max(0, len(slide.get("shapes", [])) - len(preview_shapes))
            if omitted_shapes:
                slide_preview["omitted_shapes"] = omitted_shapes
            preview_slides.append(slide_preview)
        preview["slides"] = preview_slides
        omitted_slides = max(0, len(slides) - len(preview_slides))
        if omitted_slides:
            preview["omitted_slides"] = omitted_slides

    if _json_size(preview) <= MAX_OUTPUT_CHARS:
        return _fit_reduced_json_result(preview)

    summary = {
        "success": True,
        "format": fmt,
        "file": parsed.get("file"),
        "output_format": "json",
        "truncated": True,
        "output_truncated": True,
        "source_truncated": source_truncated,
        "_warning": preview["_warning"],
    }
    for key in ("page_count", "sheet_count", "row_count", "paragraph_count", "table_count", "slide_count", "delimiter"):
        if key in parsed:
            summary[key] = parsed[key]
    if "metadata" in parsed:
        summary["metadata"] = parsed["metadata"]
    return _fit_reduced_json_result(summary)


def _fit_reduced_json_result(result: dict[str, Any]) -> dict[str, Any]:
    """Shrink reduced JSON metadata if needed to stay within the output limit."""
    if _json_size(result) <= MAX_OUTPUT_CHARS:
        return result

    if "_warning" in result:
        result["_warning"] = "Output truncated."
        if _json_size(result) <= MAX_OUTPUT_CHARS:
            return result
        result.pop("_warning", None)

    return result


def _source_data_truncated(parsed: dict[str, Any]) -> bool:
    """Return whether the parser itself truncated source data before output reduction."""
    if bool(parsed.get("truncated")):
        return True

    if parsed.get("format") == "excel":
        return any(bool(sheet.get("truncated")) for sheet in parsed.get("sheets", []))

    return False


# ---------------------------------------------------------------------------
# PDF parser
# ---------------------------------------------------------------------------


def _parse_pdf(file_path: Path) -> dict[str, Any]:
    """Extract text and metadata from a PDF file using pdfplumber."""
    try:
        import pdfplumber
    except Exception as exc:
        return _make_dependency_error("PDF_DEPENDENCY_ERROR", "pdfplumber", exc, file_path)

    pages: list[dict[str, Any]] = []
    metadata: dict[str, Any] = {}

    try:
        with pdfplumber.open(file_path) as pdf:
            metadata = {
                "page_count": len(pdf.pages),
                "metadata": {k: str(v) for k, v in (pdf.metadata or {}).items()},
            }
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                tables = page.extract_tables() or []
                page_data: dict[str, Any] = {
                    "page_number": i + 1,
                    "text": text,
                }
                if tables:
                    page_data["tables"] = [
                        [[cell if cell is not None else "" for cell in row] for row in table]
                        for table in tables
                    ]
                pages.append(page_data)
    except Exception as exc:
        return _make_error("PDF_PARSE_ERROR", f"Failed to parse PDF: {exc}", file=str(file_path))

    return {
        "success": True,
        "format": "pdf",
        "file": str(file_path),
        **metadata,
        "pages": pages,
    }


# ---------------------------------------------------------------------------
# Excel parser (.xlsx / .xlsm — read-only, data_only=True, no macro execution)
# ---------------------------------------------------------------------------


def _parse_excel(file_path: Path) -> dict[str, Any]:
    """Extract data from an Excel workbook.  XLSM is opened read-only with
    data_only=True; macros are never executed."""
    try:
        import openpyxl
    except Exception as exc:
        return _make_dependency_error("EXCEL_DEPENDENCY_ERROR", "openpyxl", exc, file_path)

    try:
        # read_only=True + data_only=True: fast streaming read, no macro execution
        wb = openpyxl.load_workbook(
            file_path,
            read_only=True,
            data_only=True,
            keep_vba=False,     # discard VBA payload entirely
            keep_links=False,
        )
    except Exception as exc:
        return _make_error("EXCEL_OPEN_ERROR", f"Failed to open workbook: {exc}", file=str(file_path))

    sheets: list[dict[str, Any]] = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list[Any]] = []
            row_count = 0
            truncated = False
            for row in ws.iter_rows(values_only=True):
                if row_count >= MAX_ROWS:
                    truncated = True
                    break
                rows.append([_excel_cell_value(cell) for cell in row])
                row_count += 1
            sheets.append({
                "sheet_name": sheet_name,
                "row_count": row_count,
                "truncated": truncated,
                "rows": rows,
            })
    except Exception as exc:
        return _make_error("EXCEL_PARSE_ERROR", f"Failed to read sheets: {exc}", file=str(file_path))
    finally:
        wb.close()

    return {
        "success": True,
        "format": "excel",
        "file": str(file_path),
        "sheet_count": len(sheets),
        "sheets": sheets,
    }


def _excel_cell_value(cell: Any) -> Any:
    """Normalise an openpyxl cell value to a JSON-safe type."""
    if cell is None:
        return None
    # datetime objects → ISO string
    if hasattr(cell, "isoformat"):
        return cell.isoformat()
    # bytes → hex (rare, but defensive)
    if isinstance(cell, bytes):
        return cell.hex()
    return cell


# ---------------------------------------------------------------------------
# CSV parser
# ---------------------------------------------------------------------------


def _parse_csv(file_path: Path, delimiter: str | None = None, encoding: str = "utf-8") -> dict[str, Any]:
    """Parse a CSV file into rows."""
    try:
        raw = file_path.read_bytes()
    except Exception as exc:
        return _make_error("CSV_READ_ERROR", f"Failed to read file: {exc}", file=str(file_path))

    # Attempt to decode
    try:
        text = raw.decode(encoding)
    except UnicodeDecodeError:
        # Fall back to latin-1 which never fails
        text = raw.decode("latin-1")

    # Auto-detect delimiter if not specified
    if delimiter is None:
        try:
            sample = text[:CSV_SNIFF_BYTES]
            dialect = csv.Sniffer().sniff(sample)
            delimiter = dialect.delimiter
        except csv.Error:
            delimiter = ","

    rows: list[list[str]] = []
    truncated = False
    try:
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        for i, row in enumerate(reader):
            if i >= MAX_ROWS:
                truncated = True
                break
            rows.append(row)
    except Exception as exc:
        return _make_error("CSV_PARSE_ERROR", f"Failed to parse CSV: {exc}", file=str(file_path))

    return {
        "success": True,
        "format": "csv",
        "file": str(file_path),
        "delimiter": delimiter,
        "row_count": len(rows),
        "truncated": truncated,
        "rows": rows,
    }


# ---------------------------------------------------------------------------
# DOCX parser
# ---------------------------------------------------------------------------


def _parse_docx(file_path: Path) -> dict[str, Any]:
    """Extract paragraphs and tables from a Word .docx file."""
    try:
        import docx
    except Exception as exc:
        return _make_dependency_error("DOCX_DEPENDENCY_ERROR", "python-docx", exc, file_path)

    try:
        doc = docx.Document(str(file_path))
    except Exception as exc:
        return _make_error("DOCX_OPEN_ERROR", f"Failed to open document: {exc}", file=str(file_path))

    paragraphs: list[dict[str, Any]] = []
    for para in doc.paragraphs:
        paragraphs.append({
            "text": para.text,
            "style": para.style.name if para.style else None,
        })

    tables: list[list[list[str]]] = []
    for table in doc.tables:
        table_rows: list[list[str]] = []
        for row in table.rows:
            table_rows.append([cell.text for cell in row.cells])
        tables.append(table_rows)

    core = doc.core_properties
    metadata: dict[str, Any] = {}
    for attr in ("author", "title", "subject", "keywords", "created", "modified"):
        val = getattr(core, attr, None)
        if val is not None:
            metadata[attr] = val.isoformat() if hasattr(val, "isoformat") else str(val)

    return {
        "success": True,
        "format": "docx",
        "file": str(file_path),
        "metadata": metadata,
        "paragraph_count": len(paragraphs),
        "table_count": len(tables),
        "paragraphs": paragraphs,
        "tables": tables,
    }


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------


def _parse_pptx(file_path: Path) -> dict[str, Any]:
    """Extract slides, shapes, and text from a PowerPoint .pptx file."""
    try:
        import pptx
    except Exception as exc:
        return _make_dependency_error("PPTX_DEPENDENCY_ERROR", "python-pptx", exc, file_path)

    try:
        prs = pptx.Presentation(str(file_path))
    except Exception as exc:
        return _make_error("PPTX_OPEN_ERROR", f"Failed to open presentation: {exc}", file=str(file_path))

    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            shape_data: dict[str, Any] = {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
            }
            if shape.has_text_frame:
                shape_data["text"] = shape.text_frame.text
            if shape.has_table:
                table_rows: list[list[str]] = []
                for row in shape.table.rows:
                    table_rows.append([cell.text for cell in row.cells])
                shape_data["table"] = table_rows
            shapes.append(shape_data)
        slides.append({
            "slide_number": i + 1,
            "shapes": shapes,
        })

    return {
        "success": True,
        "format": "pptx",
        "file": str(file_path),
        "slide_count": len(slides),
        "slides": slides,
    }


# ---------------------------------------------------------------------------
# JSON → Markdown conversion
# ---------------------------------------------------------------------------


def _to_markdown(parsed: dict[str, Any]) -> str:
    """Convert a parsed-document JSON dict into a readable Markdown string."""
    if not parsed.get("success"):
        return f"**Error**: {parsed.get('error', 'Unknown error')}"

    fmt = parsed.get("format", "unknown")
    parts: list[str] = [f"# Parsed Document: {parsed.get('file', 'unknown')}\n"]

    if fmt == "pdf":
        meta = parsed.get("metadata", {})
        if meta:
            parts.append("## Metadata\n")
            for k, v in meta.items():
                parts.append(f"- **{k}**: {v}")
            parts.append("")
        for page in parsed.get("pages", []):
            parts.append(f"## Page {page['page_number']}\n")
            parts.append(page.get("text", ""))
            for ti, table in enumerate(page.get("tables", []), 1):
                parts.append(f"\n### Table {ti}\n")
                parts.append(_table_to_md(table))
            parts.append("")

    elif fmt == "excel":
        for sheet in parsed.get("sheets", []):
            parts.append(f"## Sheet: {sheet['sheet_name']}\n")
            parts.append(f"Rows: {sheet['row_count']}")
            if sheet.get("truncated"):
                parts.append(f"*(truncated at {MAX_ROWS:,} rows)*\n")
            rows = sheet.get("rows", [])
            if rows:
                parts.append(_table_to_md([[str(c) if c is not None else "" for c in r] for r in rows]))
            parts.append("")

    elif fmt == "csv":
        parts.append(f"Delimiter: `{parsed.get('delimiter', ',')}`  ")
        parts.append(f"Rows: {parsed.get('row_count', 0)}\n")
        rows = parsed.get("rows", [])
        if rows:
            parts.append(_table_to_md(rows))

    elif fmt == "docx":
        meta = parsed.get("metadata", {})
        if meta:
            parts.append("## Metadata\n")
            for k, v in meta.items():
                parts.append(f"- **{k}**: {v}")
            parts.append("")
        for para in parsed.get("paragraphs", []):
            style = para.get("style", "")
            text = para.get("text", "")
            if style and style.startswith("Heading"):
                # e.g. "Heading 1" → "## ..."
                try:
                    level = int(style.split()[-1]) + 1
                except (ValueError, IndexError):
                    level = 2
                parts.append(f"{'#' * level} {text}\n")
            else:
                parts.append(text)
        for ti, table in enumerate(parsed.get("tables", []), 1):
            parts.append(f"\n### Table {ti}\n")
            parts.append(_table_to_md(table))

    elif fmt == "pptx":
        for slide in parsed.get("slides", []):
            parts.append(f"## Slide {slide['slide_number']}\n")
            for shape in slide.get("shapes", []):
                if "text" in shape and shape["text"].strip():
                    parts.append(shape["text"])
                if "table" in shape:
                    parts.append(_table_to_md(shape["table"]))
            parts.append("")

    return "\n".join(parts)


def _table_to_md(rows: list[list[str]]) -> str:
    """Convert a list-of-lists table to a Markdown pipe table."""
    if not rows:
        return ""
    # Compute column widths for alignment
    col_count = max(len(r) for r in rows)
    # Normalise row lengths
    normalised: list[list[str]] = []
    for row in rows:
        cells = [str(c).replace("|", "\\|").replace("\n", " ") for c in row]
        while len(cells) < col_count:
            cells.append("")
        normalised.append(cells)

    lines: list[str] = []
    # Header row
    lines.append("| " + " | ".join(normalised[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
    for row in normalised[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Dispatch helper
# ---------------------------------------------------------------------------

_EXTENSION_PARSERS: dict[str, Any] = {
    ".pdf": lambda p, **_kw: _parse_pdf(p),
    ".xlsx": lambda p, **_kw: _parse_excel(p),
    ".xlsm": lambda p, **_kw: _parse_excel(p),
    ".csv": lambda p, **kw: _parse_csv(p, delimiter=kw.get("csv_delimiter"), encoding=kw.get("csv_encoding", "utf-8")),
    ".docx": lambda p, **_kw: _parse_docx(p),
    ".pptx": lambda p, **_kw: _parse_pptx(p),
}


def _dispatch_parse(
    file_path: str,
    output_format: str = "json",
    allowed_extensions: set[str] | None = None,
    **kwargs: Any,
) -> dict[str, Any]:
    """Validate, parse, and optionally convert to Markdown."""
    # Validate output format
    fmt = output_format.lower().strip()
    if fmt not in OUTPUT_FORMATS:
        return _make_error("INVALID_FORMAT", f"output_format must be one of {OUTPUT_FORMATS}, got '{output_format}'")

    # Validate file
    resolved, err = _validate_file(file_path, allowed_extensions)
    if err:
        return _make_error("FILE_ERROR", err, file=str(resolved))

    ext = resolved.suffix.lower()
    parser = _EXTENSION_PARSERS.get(ext)
    if parser is None:
        return _make_error("UNSUPPORTED_FORMAT", f"No parser for extension '{ext}'", file=str(resolved))

    _log(f"Parsing {resolved} (format={ext}, output={fmt})")
    try:
        parsed = parser(resolved, **kwargs)
    except Exception as exc:
        return _make_error("PARSER_RUNTIME_ERROR", f"Unexpected parser failure: {exc}", file=str(resolved))
    if not isinstance(parsed, dict):
        return _make_error("PARSER_RUNTIME_ERROR", "Parser returned a non-dict result", file=str(resolved))

    # Convert to Markdown if requested
    if fmt == "markdown" and parsed.get("success"):
        try:
            md_text = _to_markdown(parsed)
        except Exception as exc:
            return _make_error("MARKDOWN_RENDER_ERROR", f"Failed to render Markdown output: {exc}", file=str(resolved))
        md_text, truncated = _truncate_output(md_text)
        return {
            "success": True,
            "format": parsed.get("format"),
            "output_format": "markdown",
            "file": str(resolved),
            "content": md_text,
            "truncated": truncated,
        }

    # JSON output — truncate serialised form if too large
    if parsed.get("success"):
        serialised = json.dumps(parsed, default=str)
        if len(serialised) > MAX_OUTPUT_CHARS:
            _log(f"Output truncated: {len(serialised):,} chars > {MAX_OUTPUT_CHARS:,} limit")
            parsed = _build_truncated_json_result(parsed, len(serialised))

    parsed["output_format"] = "json"
    return parsed


# ---------------------------------------------------------------------------
# MCP Server
# ---------------------------------------------------------------------------

mcp = FastMCP(
    "document-parser",
    instructions=(
        "Document Parser MCP server — extract structured content from PDF, "
        "Excel (.xlsx/.xlsm), CSV, Word (.docx), and PowerPoint (.pptx) files. "
        "Returns parsed data as JSON or Markdown. Read-only: never writes files "
        "or executes macros."
    ),
)


@mcp.tool()
def parse_document(
    file_path: str,
    output_format: str = "json",
    csv_delimiter: str | None = None,
    csv_encoding: str = "utf-8",
) -> dict[str, Any]:
    """Parse a document file and return its structured content.

    Auto-detects format from file extension.  Supported: .pdf, .xlsx, .xlsm,
    .csv, .docx, .pptx.

    Args:
        file_path:      Absolute or relative path to the document file.
        output_format:  ``"json"`` (default) or ``"markdown"``.
        csv_delimiter:  Optional delimiter override for CSV files (auto-detected if omitted).
        csv_encoding:   Encoding for CSV files (default ``"utf-8"``).

    Returns:
        A dict with ``success``, parsed content (format-specific fields), and
        ``output_format``.  On error, returns ``success=False`` with ``error``
        and ``error_code``.
    """
    return _dispatch_parse(
        file_path,
        output_format=output_format,
        csv_delimiter=csv_delimiter,
        csv_encoding=csv_encoding,
    )


@mcp.tool()
def parse_pdf(file_path: str, output_format: str = "json") -> dict[str, Any]:
    """Parse a PDF file and extract text and tables from each page.

    Args:
        file_path:      Path to the .pdf file.
        output_format:  ``"json"`` (default) or ``"markdown"``.

    Returns:
        Parsed page content with text and any detected tables.
    """
    return _dispatch_parse(file_path, output_format=output_format, allowed_extensions={".pdf"})


@mcp.tool()
def parse_excel(file_path: str, output_format: str = "json") -> dict[str, Any]:
    """Parse an Excel workbook (.xlsx or .xlsm) and extract sheet data.

    The file is opened read-only with data_only=True.  VBA macros in .xlsm
    files are never loaded or executed.

    Args:
        file_path:      Path to the .xlsx or .xlsm file.
        output_format:  ``"json"`` (default) or ``"markdown"``.

    Returns:
        Sheet names and row data for each sheet.
    """
    return _dispatch_parse(file_path, output_format=output_format, allowed_extensions={".xlsx", ".xlsm"})


@mcp.tool()
def parse_csv(
    file_path: str,
    output_format: str = "json",
    delimiter: str | None = None,
    encoding: str = "utf-8",
) -> dict[str, Any]:
    """Parse a CSV file and extract rows.

    Args:
        file_path:      Path to the .csv file.
        output_format:  ``"json"`` (default) or ``"markdown"``.
        delimiter:      Column delimiter (auto-detected if omitted).
        encoding:       File encoding (default ``"utf-8"``).

    Returns:
        Parsed rows with auto-detected or specified delimiter.
    """
    return _dispatch_parse(
        file_path,
        output_format=output_format,
        allowed_extensions={".csv"},
        csv_delimiter=delimiter,
        csv_encoding=encoding,
    )


@mcp.tool()
def parse_docx(file_path: str, output_format: str = "json") -> dict[str, Any]:
    """Parse a Word document (.docx) and extract paragraphs and tables.

    Args:
        file_path:      Path to the .docx file.
        output_format:  ``"json"`` (default) or ``"markdown"``.

    Returns:
        Paragraphs (with style info) and tables from the document.
    """
    return _dispatch_parse(file_path, output_format=output_format, allowed_extensions={".docx"})


@mcp.tool()
def parse_pptx(file_path: str, output_format: str = "json") -> dict[str, Any]:
    """Parse a PowerPoint presentation (.pptx) and extract slide content.

    Args:
        file_path:      Path to the .pptx file.
        output_format:  ``"json"`` (default) or ``"markdown"``.

    Returns:
        Slide-by-slide content with text and tables from each shape.
    """
    return _dispatch_parse(file_path, output_format=output_format, allowed_extensions={".pptx"})


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    _log("Starting document-parser MCP server")
    mcp.run(transport="stdio")
